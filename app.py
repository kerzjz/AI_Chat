import streamlit as st
import json
import urllib.request
import urllib.error
from datetime import datetime
import uuid
import os
import io
import re

st.set_page_config(
    page_title="Kzz AI", 
    page_icon="🤖", 
    layout="wide",
    initial_sidebar_state="expanded"
)

# MDUI Material Design 3 样式
st.markdown("""
<link href="https://fonts.googleapis.com/css2?family=Roboto:wght@400;500;700&display=swap" rel="stylesheet">
<link href="https://fonts.googleapis.com/icon?family=Material+Icons" rel="stylesheet">

<style>
    :root {
        --md-sys-color-primary: #6750A4;
        --md-sys-color-on-primary: #FFFFFF;
        --md-sys-color-primary-container: #EADDFF;
        --md-sys-color-on-primary-container: #21005D;
        --md-sys-color-secondary: #625B71;
        --md-sys-color-on-secondary: #FFFFFF;
        --md-sys-color-secondary-container: #E8DEF8;
        --md-sys-color-on-secondary-container: #1D192B;
        --md-sys-color-surface: #FFFBFE;
        --md-sys-color-surface-variant: #E7E0EC;
        --md-sys-color-on-surface: #1C1B1F;
        --md-sys-color-on-surface-variant: #49454F;
        --md-sys-color-background: #FFFBFE;
        --md-sys-color-on-background: #1C1B1F;
        --md-sys-color-error: #B3261E;
        --md-sys-color-on-error: #FFFFFF;
        --md-sys-color-outline: #79747E;
        --md-sys-color-outline-variant: #CAC4D0;
        --md-sys-elevation-1: 0px 1px 2px rgba(0,0,0,0.3), 0px 1px 3px 1px rgba(0,0,0,0.15);
        --md-sys-elevation-2: 0px 1px 2px rgba(0,0,0,0.3), 0px 2px 6px 2px rgba(0,0,0,0.15);
        --md-sys-elevation-3: 0px 1px 3px rgba(0,0,0,0.3), 0px 4px 8px 3px rgba(0,0,0,0.15);
        --md-sys-shape-corner-small: 4px;
        --md-sys-shape-corner-medium: 8px;
        --md-sys-shape-corner-large: 12px;
        --md-sys-shape-corner-extra-large: 28px;
        --md-sys-typescale-font-family: 'Roboto', sans-serif;
    }

    * { font-family: var(--md-sys-typescale-font-family); }
    .stApp { background-color: var(--md-sys-color-background); }

    [data-testid="stSidebar"] {
        background-color: var(--md-sys-color-surface);
        border-right: 1px solid var(--md-sys-color-outline-variant);
        box-shadow: var(--md-sys-elevation-1);
    }

    .stButton > button {
        border-radius: var(--md-sys-shape-corner-medium);
        font-weight: 500;
        letter-spacing: 0.5px;
        text-transform: none;
        transition: all 0.2s cubic-bezier(0.4, 0, 0.2, 1);
    }

    .stButton > button[kind="primary"] {
        background-color: var(--md-sys-color-primary);
        color: var(--md-sys-color-on-primary);
        box-shadow: none;
    }

    .stButton > button[kind="primary"]:hover {
        background-color: #7F67BE;
        box-shadow: var(--md-sys-elevation-1);
    }

    .stButton > button[kind="secondary"] {
        background-color: transparent;
        color: var(--md-sys-color-primary);
        border: 1px solid var(--md-sys-color-outline);
    }

    .stButton > button[kind="secondary"]:hover {
        background-color: rgba(103, 80, 164, 0.08);
    }

    .message-card {
        padding: 16px;
        margin: 8px 0;
        border-radius: var(--md-sys-shape-corner-large);
        box-shadow: var(--md-sys-elevation-1);
        animation: message-enter 0.3s ease-out;
    }

    @keyframes message-enter {
        from { opacity: 0; transform: translateY(8px); }
        to { opacity: 1; transform: translateY(0); }
    }

    .user-message {
        background-color: var(--md-sys-color-primary-container);
        color: var(--md-sys-color-on-primary-container);
        margin-left: 15%;
        border-bottom-right-radius: var(--md-sys-shape-corner-small);
    }

    .ai-message {
        background-color: var(--md-sys-color-surface);
        color: var(--md-sys-color-on-surface);
        margin-right: 15%;
        border: 1px solid var(--md-sys-color-outline-variant);
        border-bottom-left-radius: var(--md-sys-shape-corner-small);
    }

    .message-timestamp {
        font-size: 11px;
        font-weight: 400;
        opacity: 0.7;
        margin-top: 8px;
        display: flex;
        align-items: center;
        gap: 4px;
    }

    .user-message .message-timestamp {
        color: var(--md-sys-color-on-primary-container);
    }

    .ai-message .message-timestamp {
        color: var(--md-sys-color-on-surface-variant);
    }

    .material-icons {
        font-family: 'Material Icons';
        font-weight: normal;
        font-style: normal;
        font-size: 20px;
        line-height: 1;
        letter-spacing: normal;
        text-transform: none;
        white-space: nowrap;
        word-wrap: normal;
        direction: ltr;
        -webkit-font-feature-settings: 'liga';
        -webkit-font-smoothing: antialiased;
    }

    .app-bar {
        background-color: var(--md-sys-color-surface);
        color: var(--md-sys-color-on-surface);
        padding: 16px 24px;
        box-shadow: var(--md-sys-elevation-2);
        border-bottom: 1px solid var(--md-sys-color-outline-variant);
    }

    .chip {
        display: inline-flex;
        align-items: center;
        padding: 6px 12px;
        border-radius: var(--md-sys-shape-corner-small);
        background-color: var(--md-sys-color-surface-variant);
        color: var(--md-sys-color-on-surface-variant);
        font-size: 12px;
        font-weight: 500;
        gap: 4px;
    }

    .divider {
        height: 1px;
        background-color: var(--md-sys-color-outline-variant);
        margin: 16px 0;
    }

    .file-card {
        background-color: var(--md-sys-color-secondary-container);
        color: var(--md-sys-color-on-secondary-container);
        padding: 12px 16px;
        border-radius: var(--md-sys-shape-corner-medium);
        margin: 8px 0;
        display: flex;
        align-items: center;
        gap: 12px;
    }

    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    header {visibility: hidden;}
</style>
""", unsafe_allow_html=True)

# 配置
TOKEN = st.secrets.get("TOKEN", os.getenv("TOKEN", ""))
ACCOUNT_ID = st.secrets.get("ACCOUNT_ID", os.getenv("ACCOUNT_ID", ""))

if not TOKEN or not ACCOUNT_ID:
    st.error("⚠️ 请在 Secrets 中配置 TOKEN 和 ACCOUNT_ID")
    st.stop()

MODELS = {
    "Kimi K2.5": "@cf/moonshotai/kimi-k2.5",
    "GPT-OSS 120B": "@cf/openai/gpt-oss-120b",
    "GPT-OSS 20B": "@cf/openai/gpt-oss-20b",
    "DeepSeek R1": "@cf/deepseek-ai/deepseek-r1-distill-qwen-32b",
    "Qwen2.5 Coder": "@cf/qwen/qwen2.5-coder-32b-instruct",
    "Qwen3 30B": "@cf/qwen/qwen3-30b-a3b-fp8",
    "Gemma 3 12B": "@cf/google/gemma-3-12b-it",
    "GLM-4.7 Flash": "@cf/zai-org/glm-4.7-flash",
}

def needs_system_role(model_id):
    lower = model_id.lower()
    no_system = ["deepseek", "qwen3", "qwq", "mistral", "llama"]
    return not any(x in lower for x in no_system)

# 文件处理
SUPPORTED_EXTENSIONS = {
    '.py': 'Python', '.js': 'JavaScript', '.ts': 'TypeScript',
    '.html': 'HTML', '.htm': 'HTML', '.css': 'CSS',
    '.java': 'Java', '.cpp': 'C++', '.c': 'C', '.h': 'C/C++ Header',
    '.hpp': 'C++ Header', '.go': 'Go', '.rs': 'Rust',
    '.php': 'PHP', '.swift': 'Swift', '.kt': 'Kotlin',
    '.rb': 'Ruby', '.scala': 'Scala', '.r': 'R',
    '.m': 'MATLAB', '.sh': 'Shell', '.bash': 'Bash',
    '.ps1': 'PowerShell', '.sql': 'SQL', '.vue': 'Vue',
    '.jsx': 'JSX', '.tsx': 'TSX', '.svelte': 'Svelte',
    '.json': 'JSON', '.xml': 'XML', '.yaml': 'YAML', '.yml': 'YAML',
    '.toml': 'TOML', '.ini': 'INI', '.conf': 'Config',
    '.csv': 'CSV', '.tsv': 'TSV', '.log': 'Log',
    '.md': 'Markdown', '.markdown': 'Markdown',
    '.txt': 'Text', '.rtf': 'RTF', '.tex': 'LaTeX',
    '.docx': 'Word', '.xlsx': 'Excel', '.pptx': 'PowerPoint',
    '.dockerfile': 'Dockerfile', '.gitignore': 'Git Ignore',
    '.cmake': 'CMake', '.makefile': 'Makefile',
}

def extract_text_from_office(file_obj, ext):
    if ext == '.docx':
        try:
            import docx
            doc = docx.Document(file_obj)
            return '
'.join([para.text for para in doc.paragraphs if para.text]), None
        except ImportError:
            return None, "缺少 python-docx"
        except Exception as e:
            return None, f"docx 错误: {str(e)}"

    elif ext == '.xlsx':
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_obj)
            texts = []
            for sheet in wb.worksheets[:3]:
                texts.append(f"=== {sheet.title} ===")
                for row in sheet.iter_rows(max_row=100, values_only=True):
                    row_text = ' '.join([str(cell) for cell in row if cell])
                    if row_text.strip():
                        texts.append(row_text)
            return '
'.join(texts), None
        except ImportError:
            return None, "缺少 openpyxl"
        except Exception as e:
            return None, f"xlsx 错误: {str(e)}"

    elif ext == '.pptx':
        try:
            from pptx import Presentation
            prs = Presentation(file_obj)
            texts = []
            for i, slide in enumerate(prs.slides[:10], 1):
                texts.append(f"=== 幻灯片 {i} ===")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            return '
'.join(texts), None
        except ImportError:
            return None, "缺少 python-pptx"
        except Exception as e:
            return None, f"pptx 错误: {str(e)}"

    return None, f"不支持的格式: {ext}"

def process_file(uploaded_file):
    if uploaded_file is None:
        return None, None

    filename = uploaded_file.name
    ext = os.path.splitext(filename)[1].lower()

    if ext in ['.docx', '.xlsx', '.pptx']:
        content, error = extract_text_from_office(uploaded_file, ext)
        if error:
            return None, error
        return {
            "name": filename,
            "ext": ext,
            "type": SUPPORTED_EXTENSIONS.get(ext, "Document"),
            "content": content,
            "size": len(content)
        }, None

    try:
        content = uploaded_file.read().decode("utf-8")
    except UnicodeDecodeError:
        try:
            uploaded_file.seek(0)
            content = uploaded_file.read().decode("gbk")
        except:
            try:
                uploaded_file.seek(0)
                content = uploaded_file.read().decode("latin-1")
            except:
                return None, f"无法解码文件 {filename}"

    file_type = SUPPORTED_EXTENSIONS.get(ext, "Text")

    return {
        "name": filename,
        "ext": ext,
        "type": file_type,
        "content": content,
        "size": len(content)
    }, None

def format_file_message(file_info, user_text=""):
    name = file_info["name"]
    content = file_info["content"]
    ftype = file_info["type"]

    max_len = 4000
    if len(content) > max_len:
        content = content[:max_len] + f"
... [截断，共 {len(file_info['content'])} 字符]"

    lang = file_info["ext"].replace(".", "") if file_info["ext"] else "text"
    prompt = f"文件 `{name}` ({ftype}):
```{lang}
{content}
```"
    if user_text:
        prompt += f"

问题：{user_text}"
    else:
        prompt += "

请分析这个文件。"
    return prompt

def generate_title(text):
    if not text:
        return "新对话"

    text = text.strip()

    # 尝试提取代码中的函数名或类名
    patterns = [
        (r"def\s+(\w+)", "函数 {}"),
        (r"class\s+(\w+)", "类 {}"),
        (r"function\s+(\w+)", "函数 {}"),
    ]

    for pattern, template in patterns:
        match = re.search(pattern, text)
        if match:
            return template.format(match.group(1))

    # 尝试提取文件相关
    if "文件" in text or "`" in text:
        file_match = re.search(r"`([^`]+)`", text)
        if file_match:
            return f"文件 {file_match.group(1)[:15]}"

    # 取前15个字符
    clean_text = text.replace(chr(10), " ").strip()
    title = clean_text[:15] + ("..." if len(clean_text) > 15 else "")
    # 修复：使用原始字符串避免转义问题
    title = re.sub(r"[<>`"
]", "", title)

    return title if title else "新对话"

# Session 管理
def init_session():
    defaults = {
        "sessions": {},
        "current_id": None,
        "pending_file": None,
        "edit_mode": None,
        "show_new": False,
        "import_key": 0
    }

    for key, value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = value

    if not st.session_state.sessions:
        create_session("开始对话")

def create_session(name, model="Kimi K2.5"):
    sid = str(uuid.uuid4())[:8]
    model_id = MODELS[model]

    messages = []
    if needs_system_role(model_id):
        messages.append({
            "role": "system",
            "content": "你是智能助手，请用中文回答。",
            "time": datetime.now().isoformat()
        })

    st.session_state.sessions[sid] = {
        "id": sid,
        "name": name,
        "model": model,
        "messages": messages,
        "created_at": datetime.now().isoformat(),
        "updated_at": datetime.now().isoformat()
    }
    st.session_state.current_id = sid
    return sid

def get_session():
    cid = st.session_state.current_id
    if cid and cid in st.session_state.sessions:
        return st.session_state.sessions[cid]

    if st.session_state.sessions:
        first = list(st.session_state.sessions.keys())[0]
        st.session_state.current_id = first
        return st.session_state.sessions[first]
    return None

def delete_session(sid):
    if sid in st.session_state.sessions:
        del st.session_state.sessions[sid]
        if st.session_state.current_id == sid:
            if st.session_state.sessions:
                st.session_state.current_id = list(st.session_state.sessions.keys())[0]
            else:
                create_session("新会话")
        return True
    return False

def rename_session(sid, new_name):
    if sid in st.session_state.sessions and new_name.strip():
        st.session_state.sessions[sid]["name"] = new_name.strip()
        st.session_state.sessions[sid]["updated_at"] = datetime.now().isoformat()
        return True
    return False

def switch_model(sid, new_model):
    if sid not in st.session_state.sessions:
        return

    session = st.session_state.sessions[sid]
    old_id = MODELS[session["model"]]
    new_id = MODELS[new_model]

    old_sys = needs_system_role(old_id)
    new_sys = needs_system_role(new_id)

    session["model"] = new_model
    if old_sys != new_sys:
        if new_sys:
            session["messages"] = [{
                "role": "system",
                "content": "你是智能助手，请用中文回答。",
                "time": datetime.now().isoformat()
            }]
        else:
            session["messages"] = []
    session["updated_at"] = datetime.now().isoformat()

# API
@st.cache_data(ttl=3600)
def get_api_url(model):
    return f"https://api.cloudflare.com/client/v4/accounts/{ACCOUNT_ID}/ai/run/{model}"

def call_api(session, text, auto_title=False):
    try:
        model_id = MODELS[session["model"]]
        url = get_api_url(model_id)

        now = datetime.now().isoformat()
        session["messages"].append({
            "role": "user",
            "content": text,
            "time": now
        })

        if auto_title:
            user_msgs = [m for m in session["messages"] if m["role"] == "user"]
            if len(user_msgs) == 1 and session["name"] in ["开始对话", "新会话", "默认会话"]:
                session["name"] = generate_title(text)

        payload_msgs = [{"role": m["role"], "content": m["content"]} 
                       for m in session["messages"]]
        data = json.dumps({"messages": payload_msgs}).encode()

        req = urllib.request.Request(
            url, data=data,
            headers={"Authorization": f"Bearer {TOKEN}",
                    "Content-Type": "application/json"},
            method="POST"
        )

        with urllib.request.urlopen(req, timeout=60) as resp:
            result = json.loads(resp.read().decode())

            if result.get("success"):
                res = result.get("result", {})
                reply = res.get("response")
                if not reply and "choices" in res:
                    reply = res["choices"][0]["message"]["content"]

                if reply:
                    session["messages"].append({
                        "role": "assistant",
                        "content": reply,
                        "time": datetime.now().isoformat()
                    })
                    session["updated_at"] = datetime.now().isoformat()
                    return reply, None
                else:
                    session["messages"].pop()
                    return None, "空响应"
            else:
                session["messages"].pop()
                err = result.get("errors", [{}])[0].get("message", "未知错误")
                return None, f"API错误: {err}"

    except Exception as e:
        session["messages"].pop()
        return None, str(e)

# 导入导出
def export_session(session, fmt):
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")

    if fmt == "json":
        data = {
            "version": "2.0",
            "export_time": datetime.now().isoformat(),
            "name": session["name"],
            "model": session["model"],
            "created_at": session.get("created_at", ""),
            "messages": session["messages"]
        }
        return json.dumps(data, ensure_ascii=False, indent=2), f"{session['name']}_{ts}.json"

    elif fmt == "md":
        content = f"# {session['name']}

"
        content += f"> 模型: {session['model']}
"
        content += f"> 导出时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}

"
        content += "---

"

        for msg in session["messages"]:
            if msg["role"] == "system":
                continue
            time_str = ""
            if "time" in msg:
                try:
                    dt = datetime.fromisoformat(msg["time"])
                    time_str = dt.strftime("%m-%d %H:%M")
                except:
                    pass

            if msg["role"] == "user":
                content += f"### 👤 用户 {time_str}

{msg['content']}

"
            else:
                content += f"### 🤖 AI {time_str}

{msg['content']}

"
            content += "---

"

        return content, f"{session['name']}_{ts}.md"

    else:
        content = f"会话: {session['name']}
模型: {session['model']}

"
        for msg in session["messages"]:
            if msg["role"] != "system":
                time_str = msg.get("time", "")[:16].replace("T", " ")
                content += f"[{msg['role']}] {time_str}
{msg['content']}

"
        return content, f"{session['name']}_{ts}.txt"

def import_session(file_content):
    try:
        data = json.loads(file_content)

        def ensure_time(msg):
            if "time" not in msg:
                msg["time"] = datetime.now().isoformat()
            return msg

        if "sessions" in data:
            count = 0
            for sid, sdata in data["sessions"].items():
                new_id = str(uuid.uuid4())[:8]
                msgs = [ensure_time(m) for m in sdata.get("messages", [])]

                st.session_state.sessions[new_id] = {
                    "id": new_id,
                    "name": sdata.get("name", f"导入_{count+1}"),
                    "model": sdata.get("model", "Kimi K2.5"),
                    "messages": msgs,
                    "created_at": sdata.get("created_at", datetime.now().isoformat()),
                    "updated_at": datetime.now().isoformat()
                }
                count += 1
            return True, f"成功导入 {count} 个会话"

        sid = str(uuid.uuid4())[:8]
        msgs = [ensure_time(m) for m in data.get("messages", [])]

        st.session_state.sessions[sid] = {
            "id": sid,
            "name": data.get("name", "导入的会话"),
            "model": data.get("model", "Kimi K2.5"),
            "messages": msgs,
            "created_at": data.get("created_at", datetime.now().isoformat()),
            "updated_at": datetime.now().isoformat()
        }
        st.session_state.current_id = sid
        return True, "成功导入 1 个会话"

    except json.JSONDecodeError as e:
        return False, f"JSON 格式错误: {str(e)}"
    except Exception as e:
        return False, f"导入失败: {str(e)}"

# UI
init_session()
session = get_session()

with st.sidebar:
    st.markdown("""
    <div style="padding: 16px; text-align: center;">
        <span class="material-icons" style="font-size: 48px; color: #6750A4;">smart_toy</span>
        <h2 style="color: #1C1B1F; margin: 8px 0 0 0; font-weight: 400;">Kzz AI</h2>
        <p style="color: #49454F; font-size: 12px; margin: 4px 0 0 0;">智能对话助手</p>
    </div>
    <div class="divider"></div>
    """, unsafe_allow_html=True)

    if st.button("➕ 新建对话", use_container_width=True, type="primary"):
        st.session_state.show_new = True
        st.rerun()

    if st.session_state.show_new:
        with st.form("new_session_form"):
            name = st.text_input("会话名称（可选）", placeholder="留空自动生成标题")
            cols = st.columns([1, 1])
            with cols[0]:
                if st.form_submit_button("✓ 创建", use_container_width=True):
                    create_session(name if name else "新会话")
                    st.session_state.show_new = False
                    st.rerun()
            with cols[1]:
                if st.form_submit_button("✗ 取消", use_container_width=True):
                    st.session_state.show_new = False
                    st.rerun()

    st.markdown("<div class='divider'></div>", unsafe_allow_html=True)

    with st.expander("📎 上传文件"):
        st.caption("支持: 代码、文本、Office(docx/xlsx/pptx)")

        uploaded = st.file_uploader(
            "选择文件",
            type=list(SUPPORTED_EXTENSIONS.keys())[1:],
            label_visibility="collapsed",
            key=f"file_uploader_{st.session_state.import_key}"
        )

        if uploaded:
            with st.spinner("读取中..."):
                info, error = process_file(uploaded)

            if error:
                st.error(f"❌ {error}")
            else:
                st.session_state.pending_file = info
                st.success(f"✅ {info['name']} ({info['type']}, {info['size']} 字符)")

                if st.button("❌ 清除文件", key="clear_pending"):
                    st.session_state.pending_file = None
                    st.rerun()

    st.markdown("<div class='divider'></div>", unsafe_allow_html=True)

    st.markdown(f"<p style='color: #49454F; font-size: 12px; font-weight: 500; margin: 8px 16px;'>📂 会话列表 ({len(st.session_state.sessions)})</p>", unsafe_allow_html=True)

    sessions_sorted = sorted(
        st.session_state.sessions.values(),
        key=lambda x: x["updated_at"],
        reverse=True
    )

    for idx, s in enumerate(sessions_sorted):
        is_active = s["id"] == st.session_state.current_id
        msg_count = len([m for m in s["messages"] if m["role"] != "system"])

        with st.container():
            cols = st.columns([10, 1, 1])

            with cols[0]:
                btn_type = "primary" if is_active else "secondary"
                label = f"{s['name'][:12]}{'...' if len(s['name']) > 12 else ''} ({msg_count})"

                if st.button(
                    label,
                    key=f"sess_{s['id']}_{idx}",
                    type=btn_type,
                    use_container_width=True,
                    help=f"创建于: {s.get('created_at', '未知')[:10]}"
                ):
                    st.session_state.current_id = s["id"]
                    st.rerun()

            with cols[1]:
                if st.button(
                    "✏️",
                    key=f"rename_{s['id']}_{idx}",
                    help="重命名会话"
                ):
                    st.session_state.edit_mode = s["id"]
                    st.rerun()

            with cols[2]:
                can_delete = len(st.session_state.sessions) > 1
                if st.button(
                    "🗑️",
                    key=f"delete_{s['id']}_{idx}",
                    help="删除会话" if can_delete else "至少保留一个会话",
                    disabled=not can_delete
                ):
                    if can_delete:
                        delete_session(s["id"])
                        st.rerun()

            if st.session_state.edit_mode == s["id"]:
                with st.form(f"rename_form_{s['id']}"):
                    new_name = st.text_input(
                        "新名称",
                        value=s["name"],
                        label_visibility="collapsed"
                    )
                    c1, c2 = st.columns([1, 1])
                    with c1:
                        if st.form_submit_button("保存", use_container_width=True):
                            rename_session(s["id"], new_name)
                            st.session_state.edit_mode = None
                            st.rerun()
                    with c2:
                        if st.form_submit_button("取消", use_container_width=True):
                            st.session_state.edit_mode = None
                            st.rerun()

    st.markdown("<div class='divider'></div>", unsafe_allow_html=True)

    with st.expander("💾 导入/导出"):
        st.markdown("**📥 导入会话**")

        import_key = f"import_uploader_{st.session_state.import_key}"
        import_file = st.file_uploader(
            "选择 JSON 文件",
            type=["json"],
            key=import_key,
            label_visibility="collapsed"
        )

        if import_file is not None:
            try:
                content = import_file.read().decode("utf-8")
                success, msg = import_session(content)

                if success:
                    st.success(f"✅ {msg}")
                    st.session_state.import_key += 1
                    st.rerun()
                else:
                    st.error(f"❌ {msg}")
                    st.session_state.import_key += 1
            except Exception as e:
                st.error(f"❌ 读取失败: {str(e)}")
                st.session_state.import_key += 1

        st.markdown("<div style='height: 8px;'></div>", unsafe_allow_html=True)

        if session:
            st.markdown("**📤 导出当前**")
            fmt = st.selectbox("格式", ["JSON", "Markdown", "TXT"], key="export_fmt")
            content, fname = export_session(session, fmt.lower())

            mime = "application/json" if fmt == "JSON" else "text/plain"
            st.download_button(
                "⬇️ 下载",
                content,
                fname,
                mime=mime,
                use_container_width=True
            )

    total_msgs = sum(
        len([m for m in s["messages"] if m["role"] != "system"])
        for s in st.session_state.sessions.values()
    )
    st.markdown(
        f"<p style='color: #79747E; font-size: 11px; text-align: center; margin-top: 16px;'>"
        f"💡 {len(st.session_state.sessions)} 个会话 · {total_msgs} 条消息"
        f"</p>",
        unsafe_allow_html=True
    )

if session:
    st.markdown(f"""
    <div class="app-bar">
        <div style="display: flex; justify-content: space-between; align-items: center;">
            <h3 style="margin: 0; font-weight: 400; color: #1C1B1F;">
                💬 {session['name']}
            </h3>
            <span class="chip">
                <span class="material-icons" style="font-size: 16px;">memory</span>
                {len([m for m in session['messages'] if m['role']!='system'])} 条消息
            </span>
        </div>
    </div>
    """, unsafe_allow_html=True)

    cols = st.columns([3, 1])
    with cols[0]:
        current = session["model"]
        options = list(MODELS.keys())
        selected = st.selectbox(
            "选择模型",
            options,
            index=options.index(current),
            label_visibility="collapsed"
        )
        if selected != current:
            switch_model(session["id"], selected)
            st.rerun()

    with cols[1]:
        model_id = MODELS[session["model"]]
        st.caption(f"📦 {model_id.split('/')[-1][:20]}")

    if st.session_state.pending_file:
        f = st.session_state.pending_file
        st.markdown(f"""
        <div class="file-card">
            <span class="material-icons">description</span>
            <div>
                <div style="font-weight: 500;">{f['name']}</div>
                <div style="font-size: 12px; opacity: 0.8;">{f['type']} · {f['size']} 字符</div>
            </div>
        </div>
        """, unsafe_allow_html=True)

    for msg in session["messages"]:
        if msg["role"] == "system":
            continue

        is_user = msg["role"] == "user"
        css_class = "user-message" if is_user else "ai-message"
        icon = "👤" if is_user else "🤖"

        time_str = ""
        if "time" in msg:
            try:
                dt = datetime.fromisoformat(msg["time"])
                time_str = dt.strftime("%m-%d %H:%M")
            except:
                pass

        with st.chat_message(msg["role"]):
            st.markdown(f"""
            <div class="message-card {css_class}">
                <div style="display: flex; align-items: center; gap: 8px; margin-bottom: 8px; font-weight: 500;">
                    {icon} {"你" if is_user else "AI"}
                </div>
                {msg['content']}
                <div class="message-timestamp">
                    <span class="material-icons" style="font-size: 12px;">schedule</span>
                    {time_str}
                </div>
            </div>
            """, unsafe_allow_html=True)

    placeholder = "输入消息..."
    if st.session_state.pending_file:
        placeholder = "输入关于文件的问题（可选）..."

    if prompt := st.chat_input(placeholder):
        if st.session_state.pending_file:
            full_prompt = format_file_message(st.session_state.pending_file, prompt)
            st.session_state.pending_file = None
            display_prompt = f"📎 [带文件] {prompt}"
        else:
            full_prompt = prompt
            display_prompt = prompt

        with st.chat_message("user"):
            st.markdown(f"""
            <div class="message-card user-message">
                <div style="display: flex; align-items: center; gap: 8px; margin-bottom: 8px; font-weight: 500;">
                    👤 你
                </div>
                {display_prompt}
                <div class="message-timestamp">
                    <span class="material-icons" style="font-size: 12px;">schedule</span>
                    {datetime.now().strftime("%m-%d %H:%M")}
                </div>
            </div>
            """, unsafe_allow_html=True)

        with st.chat_message("assistant"):
            with st.spinner("思考中..."):
                reply, error = call_api(session, full_prompt, auto_title=True)

                if error:
                    st.error(f"❌ {error}")
                else:
                    st.markdown(f"""
                    <div class="message-card ai-message">
                        <div style="display: flex; align-items: center; gap: 8px; margin-bottom: 8px; font-weight: 500;">
                            🤖 AI
                        </div>
                        {reply}
                        <div class="message-timestamp">
                            <span class="material-icons" style="font-size: 12px;">schedule</span>
                            {datetime.now().strftime("%m-%d %H:%M")}
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
                    st.rerun()

    st.markdown("<div class='divider'></div>", unsafe_allow_html=True)

    col1, col2, col3 = st.columns([2, 2, 4])

    with col1:
        if st.button("🗑️ 清空对话", use_container_width=True, type="secondary"):
            model_id = MODELS[session["model"]]
            if needs_system_role(model_id):
                session["messages"] = [{
                    "role": "system",
                    "content": "你是智能助手，请用中文回答。",
                    "time": datetime.now().isoformat()
                }]
            else:
                session["messages"] = []
            session["updated_at"] = datetime.now().isoformat()
            st.session_state.pending_file = None
            st.rerun()

    with col2:
        count = len([m for m in session["messages"] if m["role"] != "system"])
        st.caption(f"📊 {count} 条消息 · {session['model']}")

else:
    st.error("⚠️ 没有活动会话")
